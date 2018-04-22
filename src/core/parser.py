import os
import urllib.request
import docx2txt
import pptx
import xlrd
from subprocess import Popen, PIPE
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
try:
    from tesserocr import PyTessBaseAPI
except:
    PyTessBaseAPI = None
from io import StringIO

class Parser:

    def __init__(self, vk_url, ready_to_use = False):
        html_code = urllib.request.urlopen(vk_url).read()
        if ready_to_use:
            self.document_text = html_code
        else:
            html_code = str(html_code).split("var src = \\'")[1]
            html_code = html_code.split("\\';")[0]
            self.document_text = urllib.request.urlopen(html_code).read()


class docParser(Parser):

    def __init__(self, vk_url):
        super().__init__(vk_url)

    def getText(self):
        file = open("test.docx", "wb")
        file.write(self.document_text)
        file.close()
        text = docx2txt.process("test.docx")
        os.remove("test.docx")
        text = "".join(text.split("\n"))
        return text


class pdfParser(Parser):

    def __init__(self, vk_url):
        super().__init__(vk_url, True)

    def getText(self):
        file = open("test.pdf", "wb")
        file.write(self.document_text)
        file.close()
        rsrcmgr = PDFResourceManager()
        retstr = StringIO()
        codec = 'utf-8'
        laparams = LAParams()
        device = TextConverter(rsrcmgr, retstr, codec=codec, laparams=laparams)
        fp = open("test.pdf", 'rb')
        interpreter = PDFPageInterpreter(rsrcmgr, device)
        password = ""
        maxpages = 0
        caching = True
        pagenos = set()

        for page in PDFPage.get_pages(fp, pagenos, maxpages=maxpages, password=password, caching=caching,
                                      check_extractable=True):
            interpreter.process_page(page)

        text = retstr.getvalue()

        fp.close()
        device.close()
        retstr.close()
        text = "".join(text.split("\n"))
        os.remove("test.pdf")
        return text


class pptxParser(Parser):

    def __init__(self, vk_url):
        super().__init__(vk_url)

    def getText(self):
        file = open("test.pptx", "wb")
        file.write(self.document_text)
        file.close()
        prs = pptx.Presentation(os.path.abspath('test.pptx'))

        text_runs = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        text = "".join(text_runs)
        os.remove("test.pdf")
        os.remove("test.pptx")
        return text


class xlsxParser(Parser):

    def __init__(self, vk_url):
        super().__init__(vk_url)

    def getText(self):
        file = open("test.xlsx", "wb")
        file.write(self.document_text)
        file.close()
        text = ""
        excel_workbook = xlrd.open_workbook("test.xlsx")

        sheet0 = excel_workbook.sheet_by_index(0)

        for row_idx in range(sheet0.nrows):
            row = sheet0.row_values(row_idx)
            text += ",".join(str(value) for value in row)
            text += "\n"
        os.remove("test.xlsx")
        return text


class txtParser(Parser):

    def __init__(self, vk_url):
        super().__init__(vk_url, True)

    def getText(self):
        file = open("test.txt", "wb")
        file.write(self.document_text)
        file.close()
        file = open("test.txt", "r")
        text = file.read()
        file.close()
        os.remove("test.txt")
        return text


class imageParser(Parser):

    def __init__(self, vk_url):
        super().__init__(vk_url, True)

    def getText(self):
        file = open("test.png", "wb")
        file.write(self.document_text)
        file.close()
        text = ""
        out, err = Popen('python ../../models/tutorials/image/imagenet/classify_image.py --image_file test.png', shell=True, stdout=PIPE).communicate()
        text += out.decode("utf-8")
        api = PyTessBaseAPI()
        api.SetImageFile("test.png")
        text += api.GetUTF8Text()
        os.remove("test.png")
        return text


ListOfExtensionsYouCanOpenWithNanoOrVim = (".txt", ".csv", ".sln", ".csproj", ".cs", ".py", ".cpp", ".c", ".hpp",
                                           ".h", ".js", ".html", ".data", ".css", ".f90", ".f")
ImageExtension = (".gif", ".jpg", ".jpeg", ".png")


def GetText(url, extension):
    base, query = url.split("?")
    query = query.split("&")
    for i in range(len(query)):
        if "no_preview" in query[i]:
            query.pop(i)
            break
    query = "&".join(query)
    url = base + "?" + query
    try:
        if extension.lower() in ListOfExtensionsYouCanOpenWithNanoOrVim:
            parser = txtParser(url)
            return parser.getText()[:32743]
        elif extension.lower() == ".pdf":
            parser = pdfParser(url)
            return parser.getText()[:32743]
        elif extension.lower() == ".doc" or extension.lower() == ".docx":
            parser = docParser(url)
            return parser.getText()[:32743]
        elif extension.lower() == ".xls" or extension.lower() == ".xlsx":
            parser = xlsxParser(url)
            return parser.getText()[:32743]
        elif extension.lower() == ".ppt" or extension.lower() == ".pptx":
            parser = pptxParser(url)
            return parser.getText()[:32743]
        elif extension.lower() in (".doc", ".docx"):
            parser = docParser(url)
            return parser.getText()[:32743]
        elif extension.lower() in (".xls", ".xlsx"):
            parser = xlsxParser(url)
            return parser.getText()[:32743]
        elif extension.lower() in (".ppt", ".pptx"):
            parser = pptxParser(url)
            return parser.getText()[:32743]
        elif extension.lower() in ImageExtension:
            parser = imageParser(url)
            return parser.getText()[:32743]
        else:
            return ""
    except:
        return ""